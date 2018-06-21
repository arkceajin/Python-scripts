from pptx import Presentation
from googletrans import Translator
import time

fileName = "07_カップ機、設置作業手順.pptx"
srcLang = 'ja'
dstLang = 'zh-CN'

prs = Presentation(fileName)
translator = Translator()

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                while True:
                    try:
                        translated = translator.translate(run.text, src=srcLang, dest=dstLang)
                        print(translated)
                    except Exception as error:
                        print('JSONDecodeError:', error, run.text)
                        time.sleep(5)
                        pass
                    else:
                        run.text = translated.text
                        break

print('Done')
prs.save('TR-' + fileName)